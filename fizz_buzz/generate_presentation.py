#!/usr/bin/env python3
"""Generate PowerPoint presentation from FizzBuzz kata git history."""

from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess


def run_git_command(cmd):
    """Run a git command and return output."""
    result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
    return result.stdout.strip()


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, content, bullet_points=None):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    if content:
        p = body.paragraphs[0]
        p.text = content
        p.level = 0

    if bullet_points:
        for i, point in enumerate(bullet_points):
            if i == 0 and not content:
                p = body.paragraphs[0]
            else:
                p = body.add_paragraph()

            if isinstance(point, tuple):
                p.text = point[0]
                p.level = point[1]
            else:
                p.text = point
                p.level = 0

    return slide


def add_code_slide(prs, title, code, language="python"):
    """Add a slide with code."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Add text box for code
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = code
    p.font.name = 'Courier New'
    p.font.size = Pt(11)

    return slide


def get_commit_details(commit_hash):
    """Get detailed information about a commit."""
    message = run_git_command(f"git log -1 --pretty=format:%B {commit_hash}")
    files = run_git_command(f"git show --name-only --pretty=format: {commit_hash}").split('\n')
    files = [f for f in files if f.strip()]
    diff = run_git_command(f"git show {commit_hash} --pretty=format:")

    return {
        'message': message,
        'files': files,
        'diff': diff
    }


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "FizzBuzz Kata with TDD",
                    "A Step-by-Step Journey Through Test-Driven Development")

    # Slide 2: What is TDD?
    add_content_slide(prs,
                          "What is Test-Driven Development?",
                      "",
                      [
                          "Red-Green-Refactor Cycle",
                          ("Write a failing test first (Red)", 1),
                          ("Write minimal code to pass the test (Green)", 1),
                          ("Refactor to improve code quality (Refactor)", 1),
                          ("Commit after each complete phase", 1),
                          "",
                          "Benefits of TDD",
                          ("Ensures code correctness from the start", 1),
                          ("Creates living documentation through tests", 1),
                          ("Makes refactoring safer", 1),
                          ("Encourages simple, focused design", 1)
                      ])

    # Slide 3: FizzBuzz Problem Statement
    add_content_slide(prs,
                          "FizzBuzz Problem",
                      "",
                      [
                          "Convert numbers to strings with these rules:",
                          ("Multiples of 3 → 'Fizz'", 1),
                          ("Multiples of 5 → 'Buzz'", 1),
                          ("Multiples of both 3 and 5 → 'FizzBuzz'", 1),
                          ("Everything else → the number as a string", 1),
                          "",
                          "Two APIs to build:",
                          ("fizzbuzz_of(number) → string", 1),
                          ("fizzbuzz_1_to(n=100) → list[string]", 1)
                      ])

    # Slide 4: Project Setup - Python Environment
    add_content_slide(prs,
                          "Project Setup: Python Environment",
                      "",
                      [
                          "Used pyenv for Python version management",
                          ("Python 3.10.13 selected", 1),
                          ("Allows easy switching between versions", 1),
                          "",
                          "Modern project structure:",
                          ("src/ - production code", 1),
                          ("src/tests/ - test code (tests inside src)", 1),
                          ("pyproject.toml - modern Python packaging", 1)
                      ])

    # Slide 5: Project Setup - pyproject.toml
    add_code_slide(prs,
                       "pyproject.toml Configuration",
                   """[project]
name = "fizz-buzz-kata"
version = "0.1.0"
requires-python = ">=3.10"

[project.optional-dependencies]
dev = [
    "pytest>=7.4.0",
    "pytest-cov>=4.1.0",
    "pytest-watch>=4.2.0",
    "flake8>=6.0.0",
    "pre-commit>=3.0.0",
]

[tool.pytest.ini_options]
testpaths = ["src/tests"]
addopts = ["-v", "--cov=src"]""")

    # Slide 6: Project Setup - Pre-commit Hooks
    add_content_slide(prs,
                          "Quality Gates: Pre-commit Hooks",
                      "",
                      [
                          "Automated checks before every commit:",
                          ("flake8 - Code style and quality checks", 1),
                          ("pytest - All tests must pass", 1),
                          ("Trailing whitespace removal", 1),
                          ("YAML/TOML validation", 1),
                          "",
                          "Benefits:",
                          ("Prevents bad code from being committed", 1),
                          ("Enforces consistent code style", 1),
                          ("Catches errors early", 1)
                      ])

    # Slide 7: Project Setup - GitHub Actions CI/CD
    add_content_slide(prs,
                          "Continuous Integration: GitHub Actions",
                      "",
                      [
                          "Automated testing on every push:",
                          ("Tests on Python 3.10, 3.11, 3.12", 1),
                          ("Runs flake8 for code quality", 1),
                          ("Runs pytest with coverage", 1),
                          ("Matrix strategy for multiple versions", 1),
                          "",
                          "Ensures code works across Python versions",
                          "Provides confidence for collaboration"
                      ])

    # Slide 8: Git Configuration
    add_content_slide(prs,
                          "Git Setup: Local vs Global Config",
                      "",
                      [
                          "Work vs Personal Projects:",
                          ("Global config for work (GitLab)", 1),
                          ("Local config for personal (GitHub)", 1),
                          "",
                          "Commands used:",
                          ("git config --local user.name 'TomSpencerLondon'", 1),
                          ("git config --local user.email 'tomspencerlondon@gmail.com'", 1),
                          "",
                          "GitHub CLI (gh) for repository management",
                          ("gh auth login", 1),
                          ("gh repo create", 1)
                      ])

    # Slide 9: Phase A Overview
    add_content_slide(prs,
                          "Phase A: Numbers (not multiples of 3 or 5)",
                      "",
                      [
                          "Goal: Return number as string for non-special numbers",
                          "",
                          "Test 1: 1 → '1'",
                          ("Red: Write failing test", 1),
                          ("Green: return '1' (fake it)", 1),
                          "",
                          "Test 2: 2 → '2'",
                          ("Red: Add test for 2", 1),
                          ("Green: if/else for 1 and 2", 1),
                          "",
                          "Test 3: 4 → '4'",
                          ("Red: Add test for 4", 1),
                          ("Green: Extend logic", 1),
                          ("Refactor: Replace with return str(number)", 1)
                      ])

    # Slide 10: Phase A Code
    add_code_slide(prs,
                       "Phase A: Final Code",
                   """def fizzbuzz_of(number):
    '''Convert a number to FizzBuzz string.'''
    return str(number)


# Tests
def test_1_is_string():
    assert fizzbuzz_of(1) == "1"

def test_2_is_string():
    assert fizzbuzz_of(2) == "2"

def test_4_is_string():
    assert fizzbuzz_of(4) == "4"
""")

    # Slide 11: Why Not Test 3 in Phase A?
    add_content_slide(prs,
                          "TDD Principle: One Behavior at a Time",
                      "",
                      [
                          "Why skip 3 in Phase A?",
                          ("3 introduces NEW behavior (Fizz)", 1),
                          ("Phase A focuses on 'number → string' behavior", 1),
                          ("Finish current behavior before adding new ones", 1),
                          "",
                          "This keeps each phase focused and simple",
                          "",
                          "Triangulation with 1, 2, 4:",
                          ("Multiple examples of same behavior", 1),
                          ("Forces general solution: str(number)", 1)
                      ])

    # Slide 12: Phase B Overview
    add_content_slide(prs,
                          "Phase B: Multiples of 3 → 'Fizz'",
                      "",
                      [
                          "Test 4: 3 → 'Fizz'",
                          ("Red: Write failing test", 1),
                          ("Green: if number % 3 == 0: return 'Fizz'", 1),
                          "",
                          "Test 5: 6 → 'Fizz' (triangulate)",
                          ("Red: Add test for 6", 1),
                          ("Green: No code change needed!", 1),
                          ("The % 3 rule already handles it", 1),
                          "",
                          "Triangulation confirms our solution is general"
                      ])

    # Slide 13: Phase B Code
    add_code_slide(prs,
                       "Phase B: Code Evolution",
                   """def fizzbuzz_of(number):
    '''Convert a number to FizzBuzz string.'''
    if number % 3 == 0:
        return "Fizz"
    return str(number)


# New tests
def test_3_is_fizz():
    assert fizzbuzz_of(3) == "Fizz"

def test_6_is_fizz():
    assert fizzbuzz_of(6) == "Fizz"
""")

    # Slide 14: Phase C Overview
    add_content_slide(prs,
                          "Phase C: Multiples of 5 → 'Buzz'",
                      "",
                      [
                          "Test 6: 5 → 'Buzz'",
                          ("Red: Write failing test", 1),
                          ("Green: if number % 5 == 0: return 'Buzz'", 1),
                          "",
                          "Test 7: 10 → 'Buzz' (triangulate)",
                          ("Red: Add test for 10", 1),
                          ("Green: No code change needed!", 1),
                          "",
                          "Pattern emerging: Test, implement, triangulate"
                      ])

    # Slide 15: Phase C Code
    add_code_slide(prs,
                       "Phase C: Adding Buzz Logic",
                   """def fizzbuzz_of(number):
    '''Convert a number to FizzBuzz string.'''
    if number % 3 == 0:
        return "Fizz"
    if number % 5 == 0:
        return "Buzz"
    return str(number)


# New tests
def test_5_is_buzz():
    assert fizzbuzz_of(5) == "Buzz"

def test_10_is_buzz():
    assert fizzbuzz_of(10) == "Buzz"
""")

    # Slide 16: Phase D Challenge
    add_content_slide(prs,
                          "Phase D: The Combined Case Challenge",
                      "",
                      [
                          "Test 8: 15 → 'FizzBuzz'",
                          "",
                          "Problem: 15 is a multiple of BOTH 3 and 5",
                          ("Current code would return 'Fizz' and stop", 1),
                          ("We need 'FizzBuzz' instead", 1),
                          "",
                          "Solution: Order matters!",
                          ("Check combined case (% 15) first", 1),
                          ("Then check individual cases", 1),
                          "",
                          "Key TDD insight: Tests drive the design"
                      ])

    # Slide 17: Phase D Code
    add_code_slide(prs,
                       "Phase D: Final Implementation",
                   """def fizzbuzz_of(number):
    '''Convert a number to FizzBuzz string.'''
    if number % 15 == 0:
        return "FizzBuzz"
    if number % 3 == 0:
        return "Fizz"
    if number % 5 == 0:
        return "Buzz"
    return str(number)


# New test
def test_15_is_fizzbuzz():
    assert fizzbuzz_of(15) == "FizzBuzz"
""")

    # Slide 18: Phase D - Alternative Approach
    add_code_slide(prs,
                       "Alternative: Build String Approach",
                   """def fizzbuzz_of(number):
    '''Alternative implementation.'''
    result = ""
    if number % 3 == 0:
        result += "Fizz"
    if number % 5 == 0:
        result += "Buzz"
    return result or str(number)

# This approach:
# - Checks both conditions independently
# - Builds the string progressively
# - No need to check % 15 explicitly
# - Same behavior, different structure""")

    # Slide 19: Phase E Overview
    add_content_slide(prs,
                          "Phase E: Sequence API",
                      "",
                      [
                          "Goal: Generate sequence from 1 to n",
                          "",
                          "Test 9: fizzbuzz_1_to(5)",
                          ("Red: Test expects ['1','2','Fizz','4','Buzz']", 1),
                          ("Green: Use list comprehension", 1),
                          ("Reuse fizzbuzz_of() for each number", 1),
                          "",
                          "Test 10: Default to 100",
                          ("Red: Test expects 100 items by default", 1),
                          ("Green: Add default parameter n=100", 1)
                      ])

    # Slide 20: Phase E Code
    add_code_slide(prs,
                       "Phase E: Sequence Implementation",
                   """def fizzbuzz_1_to(n=100):
    '''Generate FizzBuzz sequence from 1 to n.

    Args:
        n: Upper limit (default 100)

    Returns:
        List of FizzBuzz strings from 1 to n
    '''
    return [fizzbuzz_of(i) for i in range(1, n + 1)]


# Tests
def test_sequence_to_5():
    assert fizzbuzz_1_to(5) == ["1", "2", "Fizz", "4", "Buzz"]

def test_sequence_default_is_100():
    result = fizzbuzz_1_to()
    assert len(result) == 100
    assert result[99] == "Buzz"
""")

    # Slide 21: Code Reuse Benefits
    add_content_slide(prs,
                          "Power of Code Reuse",
                      "",
                      [
                          "fizzbuzz_1_to reuses fizzbuzz_of",
                          "",
                          "Benefits:",
                          ("Single source of truth for FizzBuzz logic", 1),
                          ("Changes to rules only need one place", 1),
                          ("Composition over duplication", 1),
                          ("Easy to test and maintain", 1),
                          "",
                          "This is the 'DRY' principle:",
                          ("Don't Repeat Yourself", 1)
                      ])

    # Slide 22: Final Test Results
    add_content_slide(prs,
                          "Final Test Results",
                      "",
                      [
                          "10 tests, all passing ✓",
                          "",
                          "Coverage: 100%",
                          ("Every line of production code is tested", 1),
                          ("High confidence in correctness", 1),
                          "",
                          "Tests run in multiple places:",
                          ("Pre-commit hooks (local)", 1),
                          ("GitHub Actions (CI/CD)", 1),
                          ("Python 3.10, 3.11, 3.12", 1),
                          "",
                          "Total time: ~30 seconds per full CI run"
                      ])

    # Slide 23: Git History
    commits = run_git_command("git log --oneline --reverse").split('\n')
    commit_bullets = []
    for commit in commits:
        if commit.strip():
            parts = commit.split(' ', 1)
            if len(parts) == 2:
                commit_bullets.append(f"{parts[1]}")

    add_content_slide(prs,
                          "Git History: The Story",
                      "",
                      commit_bullets)

    # Slide 24: Commit Strategy
    add_content_slide(prs,
                          "Commit Strategy for Learning",
                      "",
                      [
                          "Commit after each phase completion",
                          "",
                          "Each commit message includes:",
                          ("Phase identifier (A, B, C, D, E)", 1),
                          ("What tests were added", 1),
                          ("What production code changed", 1),
                          ("Test count and coverage", 1),
                          "",
                          "Benefits for revision:",
                          ("Easy to replay the progression", 1),
                          ("See how tests drive design", 1),
                          ("Understand each decision point", 1)
                      ])

    # Slide 25: Key TDD Learnings
    add_content_slide(prs,
                          "Key TDD Learnings",
                      "",
                      [
                          "1. Write ONE failing test at a time",
                          ("Focus prevents overwhelm", 1),
                          "",
                          "2. Write minimal code to pass",
                          ("'Fake it till you make it' is OK", 1),
                          "",
                          "3. Triangulation forces general solutions",
                          ("Multiple examples reveal patterns", 1),
                          "",
                          "4. Refactor only when tests are green",
                          ("Safety net prevents breaking changes", 1),
                          "",
                          "5. Test one behavior at a time",
                          ("Keeps phases focused and manageable", 1)
                      ])

    # Slide 26: Quality Automation Stack
    add_content_slide(prs,
                          "Quality Automation Stack",
                      "",
                      [
                          "Layer 1: Pre-commit Hooks (Local)",
                          ("Runs before each commit", 1),
                          ("Fast feedback (< 5 seconds)", 1),
                          ("Prevents bad commits", 1),
                          "",
                          "Layer 2: GitHub Actions (Remote)",
                          ("Runs on push to GitHub", 1),
                          ("Tests multiple Python versions", 1),
                          ("Catches environment-specific issues", 1),
                          "",
                          "Result: High confidence, low manual effort"
                      ])

    # Slide 27: Project Structure
    add_code_slide(prs,
                       "Final Project Structure",
                   """fizz_buzz/
├── .github/
│   └── workflows/
│       └── ci.yml          # GitHub Actions
├── src/
│   ├── __init__.py
│   ├── fizzbuzz.py         # Production code
│   └── tests/
│       ├── __init__.py
│       └── test_fizzbuzz.py # Tests
├── .pre-commit-config.yaml  # Hooks config
├── .flake8                  # Linting rules
├── pyproject.toml           # Project config
├── CLAUDE.md                # AI pairing guide
└── README.md                # Documentation""")

    # Slide 28: Refactoring Opportunities
    add_content_slide(prs,
                          "Potential Refactorings (Stretch Goals)",
                      "",
                      [
                          "Make divisors configurable:",
                          ("rules = [(3, 'Fizz'), (5, 'Buzz')]", 1),
                          ("More flexible for variations", 1),
                          "",
                          "Add input validation:",
                          ("Reject negative numbers", 1),
                          ("Reject non-integers", 1),
                          "",
                          "Output formatting:",
                          ("Join with newlines or spaces", 1),
                          ("Print directly to console", 1),
                          "",
                          "All these would be test-driven!"
                      ])

    # Slide 29: Lessons for Future Katas
    add_content_slide(prs,
                          "Applying These Lessons",
                      "",
                      [
                          "This approach works for any kata:",
                          "",
                          "1. Set up quality gates first",
                          ("Pre-commit hooks", 1),
                          ("CI/CD pipeline", 1),
                          "",
                          "2. Break problem into phases",
                          ("One behavior per phase", 1),
                          ("Simple to complex", 1),
                          "",
                          "3. Follow Red-Green-Refactor strictly",
                          "",
                          "4. Commit regularly with good messages",
                          "",
                          "5. Use git history as learning tool"
                      ])

    # Slide 30: Summary
    add_content_slide(prs,
                          "Summary: FizzBuzz TDD Journey",
                      "",
                      [
                          "✓ Modern Python project setup",
                          "✓ Automated quality gates (hooks + CI/CD)",
                          "✓ 5 phases, each building on the last",
                          "✓ 10 tests, 100% coverage",
                          "✓ Clean, maintainable code",
                          "✓ Git history tells the story",
                          "",
                          "TDD isn't just about testing—",
                          "it's a design methodology that leads to:",
                          ("Better architecture", 1),
                          ("Higher confidence", 1),
                          ("Living documentation", 1),
                          ("Easier refactoring", 1)
                      ])

    return prs


if __name__ == "__main__":
    print("Generating FizzBuzz TDD presentation...")
    prs = create_presentation()

    output_file = "FizzBuzz_TDD_Journey.pptx"
    prs.save(output_file)
    print(f"✓ Presentation saved as {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
