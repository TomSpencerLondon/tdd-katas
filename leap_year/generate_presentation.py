"""Generate PowerPoint presentation and PDF for Leap Year Kata."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.colors import HexColor


def create_powerpoint():
    """Create PowerPoint presentation for Leap Year Kata."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors
    TITLE_COLOR = RGBColor(31, 78, 121)
    CODE_BG = RGBColor(43, 43, 43)
    CODE_TEXT = RGBColor(212, 212, 212)
    GREEN = RGBColor(76, 175, 80)
    RED = RGBColor(244, 67, 54)

    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide

    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.space_before = Pt(12)
            p.level = 0
        
        return slide

    def add_code_slide(title, code, phase_info=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        
        # Phase info
        if phase_info:
            info_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.4))
            info_frame = info_box.text_frame
            info_frame.text = phase_info
            info_frame.paragraphs[0].font.size = Pt(16)
            info_frame.paragraphs[0].font.italic = True
            code_top = 1.8
        else:
            code_top = 1.5
        
        # Code box
        code_box = slide.shapes.add_textbox(Inches(0.8), Inches(code_top), Inches(8.4), Inches(5))
        code_frame = code_box.text_frame
        code_frame.word_wrap = True
        
        for line in code.split('\n'):
            p = code_frame.add_paragraph()
            p.text = line
            p.font.name = 'Courier New'
            p.font.size = Pt(14)
            p.font.color.rgb = CODE_TEXT
        
        # Add background
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.7), Inches(code_top - 0.1),
            Inches(8.6), Inches(5.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CODE_BG
        shape.line.color.rgb = CODE_BG
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)
        
        return slide

    def add_tdd_cycle_slide(phase, test_desc, result, code_change, learning):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"Phase {phase}: {test_desc}"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
        
        # RED section
        red_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(1))
        red_frame = red_box.text_frame
        p = red_frame.paragraphs[0]
        p.text = f"🔴 RED: {result}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RED
        
        # GREEN section
        green_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(8.4), Inches(2.5))
        green_frame = green_box.text_frame
        p = green_frame.paragraphs[0]
        p.text = "🟢 GREEN: Production Code"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = GREEN
        
        for line in code_change.split('\n'):
            p = green_frame.add_paragraph()
            p.text = line
            p.font.name = 'Courier New'
            p.font.size = Pt(14)
            p.space_before = Pt(4)
        
        # LEARNING section
        learn_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(1.8))
        learn_frame = learn_box.text_frame
        p = learn_frame.paragraphs[0]
        p.text = f"💡 Learning: {learning}"
        p.font.size = Pt(16)
        p.font.italic = True
        
        return slide

    # Slide 1: Title
    add_title_slide(
        "Leap Year Kata",
        "A Test-Driven Development Journey"
    )

    # Slide 2: What is a Leap Year?
    add_content_slide(
        "What is a Leap Year?",
        [
            "• A year with an extra day (February 29th)",
            "• Aligns the calendar year with the solar year",
            "• Before 1582: Julian Calendar (every 4 years)",
            "• After 1582: Gregorian Calendar (more accurate)",
            "",
            "Goal: Implement is_leap_year(year) → bool"
        ]
    )

    # Slide 3: The Rules
    add_content_slide(
        "Gregorian Calendar Rules",
        [
            "1️⃣  Years divisible by 400 ARE leap years",
            "     Example: 2000 ✓",
            "",
            "2️⃣  Years divisible by 100 but not 400 are NOT leap years",
            "     Example: 1900 ✗",
            "",
            "3️⃣  Years divisible by 4 but not 100 ARE leap years",
            "     Example: 2016 ✓",
            "",
            "4️⃣  Years not divisible by 4 are NOT leap years",
            "     Example: 2017 ✗"
        ]
    )

    # Slide 4: TDD Approach
    add_content_slide(
        "Our TDD Approach",
        [
            "🔴 RED: Write a failing test",
            "",
            "🟢 GREEN: Write the simplest code to pass",
            "",
            "🔵 REFACTOR: Improve the code",
            "",
            "📝 COMMIT: Document the progress",
            "",
            "🔁 REPEAT: One test at a time"
        ]
    )

    # Slide 5: The Four Phases
    add_content_slide(
        "The Four Phases",
        [
            "Phase A: Non-leap years (2017, 2018, 2019)",
            "         → Start simple: return False",
            "",
            "Phase B: Basic leap years (2016, 2012, 2008)",
            "         → Add: if year % 4 == 0",
            "",
            "Phase C: Century exceptions (1900, 1800, 1700, 2100)",
            "         → Add: if year % 100 == 0",
            "",
            "Phase D: 400-year rule (2000, 2400, 1600)",
            "         → Add: if year % 400 == 0"
        ]
    )

    # Phase A slides
    add_tdd_cycle_slide(
        "A",
        "Test 2017 → False",
        "Test fails (returns None)",
        "def is_leap_year(year):\n    return False",
        "\"Fake it till you make it\" - start with hardcoded values"
    )

    add_code_slide(
        "Phase A: Complete",
        """def test_2017_is_not_leap_year():
    assert is_leap_year(2017) is False

def test_2018_is_not_leap_year():
    assert is_leap_year(2018) is False

def test_2019_is_not_leap_year():
    assert is_leap_year(2019) is False""",
        "✅ 3 tests passing | Production: return False"
    )

    # Phase B slides
    add_tdd_cycle_slide(
        "B",
        "Test 2016 → True",
        "Test fails (returns False, expected True)",
        "def is_leap_year(year):\n    if year % 4 == 0:\n        return True\n    return False",
        "Tests force us to implement real logic, not hardcoded values"
    )

    add_code_slide(
        "Phase B: Complete",
        """def test_2016_is_leap_year():
    assert is_leap_year(2016) is True

def test_2012_is_leap_year():
    assert is_leap_year(2012) is True

def test_2008_is_leap_year():
    assert is_leap_year(2008) is True""",
        "✅ 6 tests passing | Triangulation confirms modulo 4 logic"
    )

    # Phase C slides
    add_tdd_cycle_slide(
        "C",
        "Test 1900 → False",
        "Test fails (returns True, expected False)",
        "if year % 100 == 0 and year % 400 != 0:\n    return False\nif year % 4 == 0:\n    return True\nreturn False",
        "Order matters! Check exceptions BEFORE general rules"
    )

    add_code_slide(
        "Phase C: Complete",
        """def test_1900_is_not_leap_year():
    assert is_leap_year(1900) is False

def test_1800_is_not_leap_year():
    assert is_leap_year(1800) is False

def test_1700_is_not_leap_year():
    assert is_leap_year(1700) is False

def test_2100_is_not_leap_year():
    assert is_leap_year(2100) is False""",
        "✅ 10 tests passing | Century years handled correctly"
    )

    # Phase D slides
    add_tdd_cycle_slide(
        "D",
        "Test 2000 → True",
        "Test passes! (Logic accidentally worked)",
        "# Refactored for clarity:\nif year % 400 == 0:\n    return True\nif year % 100 == 0:\n    return False\nif year % 4 == 0:\n    return True\nreturn False",
        "Make implicit logic explicit - check most specific rules first"
    )

    add_code_slide(
        "Phase D: Complete",
        """def test_2000_is_leap_year():
    assert is_leap_year(2000) is True

def test_2400_is_leap_year():
    assert is_leap_year(2400) is True

def test_1600_is_leap_year():
    assert is_leap_year(1600) is True""",
        "✅ 13 tests passing | All rules implemented!"
    )

    # Final implementation
    add_code_slide(
        "Final Implementation",
        """def is_leap_year(year):
    \"\"\"Determine if a year is a leap year.\"\"\"
    # Years divisible by 400 ARE leap years
    if year % 400 == 0:
        return True
    # Century years (divisible by 100) are NOT leap years
    if year % 100 == 0:
        return False
    # Years divisible by 4 are leap years
    if year % 4 == 0:
        return True
    return False""",
        "✅ 13 tests | 100% coverage | All Gregorian rules implemented"
    )

    # Alternative implementation
    add_code_slide(
        "Alternative: Boolean Expression",
        """def is_leap_year(year):
    \"\"\"Condensed version using boolean logic.\"\"\"
    return (year % 4 == 0 and year % 100 != 0) or (year % 400 == 0)


# This is equivalent to:
# - Divisible by 4 AND NOT a century year, OR
# - Divisible by 400

# Both implementations are correct!
# The explicit version is more readable.
# The condensed version is more elegant.""",
        "Same logic, different style"
    )

    # Test coverage breakdown
    add_content_slide(
        "Test Coverage Breakdown",
        [
            "Phase A: Non-leap years (3 tests)",
            "  2017, 2018, 2019 → False",
            "",
            "Phase B: Basic leap years (3 tests)",
            "  2016, 2012, 2008 → True",
            "",
            "Phase C: Century exceptions (4 tests)",
            "  1900, 1800, 1700, 2100 → False",
            "",
            "Phase D: 400-year rule (3 tests)",
            "  2000, 2400, 1600 → True",
            "",
            "Total: 13 tests, 100% code coverage ✅"
        ]
    )

    # Key learnings
    add_content_slide(
        "Key TDD Principles",
        [
            "1. Red-Green-Refactor: Every change followed this cycle",
            "",
            "2. Fake it till you make it: Start simple, add complexity",
            "",
            "3. Triangulation: Multiple tests confirm general solutions",
            "",
            "4. Small steps: One test at a time",
            "",
            "5. Order matters: Most specific rules first (400→100→4)",
            "",
            "6. Frequent commits: Clear history showing progression"
        ]
    )

    # Git history
    add_content_slide(
        "Git Commit History",
        [
            "fe5f47c  Phase A: Non-leap years (not divisible by 4)",
            "",
            "23759ba  Phase B: Basic leap years (divisible by 4)",
            "",
            "2969d65  Phase C: Century years NOT divisible by 400",
            "",
            "bb08149  Phase D: Century years divisible by 400",
            "",
            "Each commit includes:",
            "  • Red-Green-Refactor details",
            "  • Test counts and coverage",
            "  • Learning points"
        ]
    )

    # What we learned
    add_content_slide(
        "What We Learned",
        [
            "✅ TDD drives design incrementally",
            "",
            "✅ Start simple, let tests demand complexity",
            "",
            "✅ Rule order matters (check 400 before 100 before 4)",
            "",
            "✅ Explicit code is better than implicit",
            "",
            "✅ Triangulation validates general solutions",
            "",
            "✅ Small commits create a reviewable story"
        ]
    )

    # Final slide
    add_title_slide(
        "Complete! 🎉",
        "13 tests passing | 100% coverage | 4 clear commits"
    )

    # Save
    prs.save('leap_year_kata_presentation.pptx')
    print("✅ PowerPoint presentation created: leap_year_kata_presentation.pptx")


def create_pdf():
    """Create PDF document for Leap Year Kata."""
    doc = SimpleDocTemplate("leap_year_kata_guide.pdf", pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Title'],
        fontSize=28,
        textColor=HexColor('#1f4e79'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading1'],
        fontSize=20,
        textColor=HexColor('#1f4e79'),
        spaceAfter=12,
        spaceBefore=16
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=HexColor('#2c5f8d'),
        spaceAfter=10,
        spaceBefore=12
    )
    
    code_style = ParagraphStyle(
        'Code',
        parent=styles['Code'],
        fontSize=10,
        fontName='Courier',
        leftIndent=20,
        spaceAfter=12,
        spaceBefore=6
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=11,
        spaceAfter=12
    )
    
    # Title
    story.append(Paragraph("Leap Year Kata", title_style))
    story.append(Paragraph("A Test-Driven Development Journey", styles['Heading2']))
    story.append(Spacer(1, 0.3*inch))
    
    # Introduction
    story.append(Paragraph("Introduction", heading_style))
    story.append(Paragraph(
        "This document walks through the complete Test-Driven Development (TDD) implementation "
        "of the Leap Year kata. We'll build the solution incrementally, one test at a time, "
        "following the Red-Green-Refactor cycle.",
        body_style
    ))
    
    # Problem statement
    story.append(Paragraph("Problem Statement", heading_style))
    story.append(Paragraph(
        "Implement a function <b>is_leap_year(year)</b> that determines if a year is a leap year "
        "according to the Gregorian Calendar rules:",
        body_style
    ))
    story.append(Paragraph("1. Years divisible by 400 ARE leap years (e.g., 2000)", body_style))
    story.append(Paragraph("2. Years divisible by 100 but not by 400 are NOT leap years (e.g., 1900)", body_style))
    story.append(Paragraph("3. Years divisible by 4 but not by 100 ARE leap years (e.g., 2016)", body_style))
    story.append(Paragraph("4. Years not divisible by 4 are NOT leap years (e.g., 2017)", body_style))
    
    # TDD Approach
    story.append(Paragraph("TDD Approach", heading_style))
    story.append(Paragraph("We follow the Red-Green-Refactor cycle:", body_style))
    story.append(Paragraph("🔴 <b>RED</b>: Write a failing test", body_style))
    story.append(Paragraph("🟢 <b>GREEN</b>: Write the simplest code to make it pass", body_style))
    story.append(Paragraph("🔵 <b>REFACTOR</b>: Improve the code without changing behavior", body_style))
    story.append(Paragraph("📝 <b>COMMIT</b>: Document the progress with a clear commit message", body_style))
    story.append(Paragraph("🔁 <b>REPEAT</b>: Continue with the next test", body_style))
    
    story.append(PageBreak())
    
    # Phase A
    story.append(Paragraph("Phase A: Non-Leap Years", heading_style))
    story.append(Paragraph("Start with the simplest case: years that are NOT leap years.", body_style))
    
    story.append(Paragraph("Test 1: 2017 → False", subheading_style))
    story.append(Paragraph("🔴 RED: Write the test", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def test_2017_is_not_leap_year():<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;assert is_leap_year(2017) is False</font>",
        code_style
    ))
    story.append(Paragraph("🟢 GREEN: Make it pass with the simplest code", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def is_leap_year(year):<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;return False</font>",
        code_style
    ))
    story.append(Paragraph(
        "💡 <b>Learning</b>: \"Fake it till you make it\" - start with hardcoded values.",
        body_style
    ))
    
    story.append(Paragraph("Tests 2 & 3: 2018, 2019 → False", subheading_style))
    story.append(Paragraph(
        "Add two more tests for triangulation. They both pass with the existing hardcoded "
        "<font name='Courier'>return False</font>. This confirms our simple solution works for "
        "all non-leap years.",
        body_style
    ))
    story.append(Paragraph("<b>Result</b>: 3 tests passing ✅", body_style))
    
    story.append(PageBreak())
    
    # Phase B
    story.append(Paragraph("Phase B: Basic Leap Years", heading_style))
    story.append(Paragraph("Now add years that ARE leap years (divisible by 4).", body_style))
    
    story.append(Paragraph("Test 4: 2016 → True", subheading_style))
    story.append(Paragraph("🔴 RED: Write the test (it will fail!)", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def test_2016_is_leap_year():<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;assert is_leap_year(2016) is True</font>",
        code_style
    ))
    story.append(Paragraph("Test fails: returns False, expected True", body_style))
    story.append(Paragraph("🟢 GREEN: Add real logic", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def is_leap_year(year):<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 4 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return True<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;return False</font>",
        code_style
    ))
    story.append(Paragraph(
        "💡 <b>Learning</b>: Tests force us to implement general logic, not hardcoded values.",
        body_style
    ))
    
    story.append(Paragraph("Tests 5 & 6: 2012, 2008 → True", subheading_style))
    story.append(Paragraph(
        "Add triangulation tests. They pass with existing logic, confirming the modulo 4 rule "
        "is general.",
        body_style
    ))
    story.append(Paragraph("<b>Result</b>: 6 tests passing ✅", body_style))
    
    story.append(PageBreak())
    
    # Phase C
    story.append(Paragraph("Phase C: Century Exceptions", heading_style))
    story.append(Paragraph("Century years (divisible by 100) are NOT leap years unless divisible by 400.", body_style))
    
    story.append(Paragraph("Test 7: 1900 → False", subheading_style))
    story.append(Paragraph("🔴 RED: Write the test (it will break our logic!)", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def test_1900_is_not_leap_year():<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;assert is_leap_year(1900) is False</font>",
        code_style
    ))
    story.append(Paragraph(
        "Test fails: returns True (because 1900 % 4 == 0), expected False",
        body_style
    ))
    story.append(Paragraph("🟢 GREEN: Add century check BEFORE the modulo 4 check", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def is_leap_year(year):<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 100 == 0 and year % 400 != 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return False<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 4 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return True<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;return False</font>",
        code_style
    ))
    story.append(Paragraph(
        "💡 <b>Learning</b>: Order matters! Check exceptions BEFORE general rules.",
        body_style
    ))
    
    story.append(Paragraph("Tests 8-10: 1800, 1700, 2100 → False", subheading_style))
    story.append(Paragraph(
        "Add more century year tests. They all pass with existing logic.",
        body_style
    ))
    story.append(Paragraph("<b>Result</b>: 10 tests passing ✅", body_style))
    
    story.append(PageBreak())
    
    # Phase D
    story.append(Paragraph("Phase D: 400-Year Rule", heading_style))
    story.append(Paragraph("Century years divisible by 400 ARE leap years.", body_style))
    
    story.append(Paragraph("Test 11: 2000 → True", subheading_style))
    story.append(Paragraph("🔴 RED: Write the test", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def test_2000_is_leap_year():<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;assert is_leap_year(2000) is True</font>",
        code_style
    ))
    story.append(Paragraph(
        "Surprisingly, test passes! Our logic accidentally worked because 2000 % 400 == 0, "
        "so the century check evaluates to False and we fall through to year % 4 == 0.",
        body_style
    ))
    story.append(Paragraph("🔵 REFACTOR: Make the 400-year rule explicit", body_style))
    story.append(Paragraph(
        "<font name='Courier'>def is_leap_year(year):<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 400 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return True<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 100 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return False<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 4 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return True<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;return False</font>",
        code_style
    ))
    story.append(Paragraph(
        "💡 <b>Learning</b>: Make implicit logic explicit. Check most specific rules first (400 → 100 → 4).",
        body_style
    ))
    
    story.append(Paragraph("Tests 12 & 13: 2400, 1600 → True", subheading_style))
    story.append(Paragraph(
        "Add triangulation tests. They pass with refactored logic.",
        body_style
    ))
    story.append(Paragraph("<b>Result</b>: 13 tests passing ✅ | 100% coverage ✅", body_style))
    
    story.append(PageBreak())
    
    # Final implementation
    story.append(Paragraph("Final Implementation", heading_style))
    story.append(Paragraph("Explicit version (more readable):", subheading_style))
    story.append(Paragraph(
        "<font name='Courier'>def is_leap_year(year):<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;\"\"\"Determine if a year is a leap year.\"\"\"<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;# Years divisible by 400 ARE leap years<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 400 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return True<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;# Century years are NOT leap years<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 100 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return False<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;# Years divisible by 4 are leap years<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;if year % 4 == 0:<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;return True<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;return False</font>",
        code_style
    ))
    
    story.append(Paragraph("Condensed version (more elegant):", subheading_style))
    story.append(Paragraph(
        "<font name='Courier'>def is_leap_year(year):<br/>"
        "&nbsp;&nbsp;&nbsp;&nbsp;return (year % 4 == 0 and year % 100 != 0) or (year % 400 == 0)</font>",
        code_style
    ))
    
    story.append(PageBreak())
    
    # Test coverage
    story.append(Paragraph("Test Coverage Breakdown", heading_style))
    story.append(Paragraph("<b>Phase A: Non-leap years (3 tests)</b>", body_style))
    story.append(Paragraph("2017, 2018, 2019 → False", body_style))
    story.append(Spacer(1, 0.1*inch))
    
    story.append(Paragraph("<b>Phase B: Basic leap years (3 tests)</b>", body_style))
    story.append(Paragraph("2016, 2012, 2008 → True", body_style))
    story.append(Spacer(1, 0.1*inch))
    
    story.append(Paragraph("<b>Phase C: Century exceptions (4 tests)</b>", body_style))
    story.append(Paragraph("1900, 1800, 1700, 2100 → False", body_style))
    story.append(Spacer(1, 0.1*inch))
    
    story.append(Paragraph("<b>Phase D: 400-year rule (3 tests)</b>", body_style))
    story.append(Paragraph("2000, 2400, 1600 → True", body_style))
    story.append(Spacer(1, 0.1*inch))
    
    story.append(Paragraph("<b>Total: 13 tests, 100% code coverage ✅</b>", body_style))
    
    # Key principles
    story.append(Paragraph("Key TDD Principles", heading_style))
    story.append(Paragraph("1. <b>Red-Green-Refactor</b>: Every change followed this cycle", body_style))
    story.append(Paragraph("2. <b>Fake it till you make it</b>: Started with hardcoded values", body_style))
    story.append(Paragraph("3. <b>Triangulation</b>: Multiple tests confirm general solutions", body_style))
    story.append(Paragraph("4. <b>Small steps</b>: One test at a time, minimal code changes", body_style))
    story.append(Paragraph("5. <b>Order matters</b>: Most specific rules checked first (400→100→4)", body_style))
    story.append(Paragraph("6. <b>Frequent commits</b>: Clear history showing progression", body_style))
    
    story.append(PageBreak())
    
    # Git history
    story.append(Paragraph("Git Commit History", heading_style))
    story.append(Paragraph(
        "<font name='Courier'>fe5f47c  Phase A: Non-leap years (not divisible by 4)</font>",
        body_style
    ))
    story.append(Paragraph(
        "<font name='Courier'>23759ba  Phase B: Basic leap years (divisible by 4)</font>",
        body_style
    ))
    story.append(Paragraph(
        "<font name='Courier'>2969d65  Phase C: Century years NOT divisible by 400</font>",
        body_style
    ))
    story.append(Paragraph(
        "<font name='Courier'>bb08149  Phase D: Century years divisible by 400</font>",
        body_style
    ))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph(
        "Each commit includes Red-Green-Refactor details, test counts, coverage, and learning points.",
        body_style
    ))
    
    # What we learned
    story.append(Paragraph("What We Learned", heading_style))
    story.append(Paragraph("✅ TDD drives design incrementally", body_style))
    story.append(Paragraph("✅ Start simple, let tests demand complexity", body_style))
    story.append(Paragraph("✅ Rule order matters (check 400 before 100 before 4)", body_style))
    story.append(Paragraph("✅ Explicit code is better than implicit", body_style))
    story.append(Paragraph("✅ Triangulation validates general solutions", body_style))
    story.append(Paragraph("✅ Small commits create a reviewable story", body_style))
    
    # Conclusion
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("Conclusion", heading_style))
    story.append(Paragraph(
        "This kata demonstrates the power of TDD. We started with no implementation, "
        "let tests drive the design, and arrived at a correct, well-tested solution with "
        "100% coverage. The journey is as valuable as the destination!",
        body_style
    ))
    
    # Build PDF
    doc.build(story)
    print("✅ PDF guide created: leap_year_kata_guide.pdf")


if __name__ == '__main__':
    print("Generating Leap Year Kata presentation materials...\n")
    create_powerpoint()
    create_pdf()
    print("\n✅ All materials generated successfully!")
    print("\nFiles created:")
    print("  • leap_year_kata_presentation.pptx")
    print("  • leap_year_kata_guide.pdf")
